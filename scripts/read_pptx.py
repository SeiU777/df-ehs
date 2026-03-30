"""讀取 PPTX 簡報內容，輸出為文字檔"""
import sys
import subprocess
import os

try:
    from pptx import Presentation
except ImportError:
    print("正在安裝 python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

pptx_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "報告", "報告.pptx")
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "報告", "報告_output.txt")

prs = Presentation(pptx_path)

with open(output_path, "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides, 1):
        f.write(f"\n{'='*60}\n")
        f.write(f"=== 第 {i} 頁投影片 ===\n")
        f.write(f"{'='*60}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        f.write(text + "\n")
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    f.write(" | ".join(cells) + "\n")
                f.write("\n")
        f.write("\n")

print(f"已輸出至: {output_path}")
print("請回到 Claude 繼續分析。")
